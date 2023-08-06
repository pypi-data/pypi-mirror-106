import os
import docx2txt
import xlrd
import csv
import cv2
import pytesseract
from PIL import Image
from pkgutil import find_loader
import PyPDF2
import pdfplumber
from pptx import Presentation
from pdf2image import convert_from_path
import speech_recognition as sr

from pyostie.convert import *
from pyostie.insights_ext import *

pandas_installed = find_loader("pandas") is not None
if pandas_installed:
    import pandas as pd

a = pd.DataFrame()
ocr_dict_output = []


class DOCXParser:

    def __init__(self, filename, img_dir):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename
        self.img_dir = img_dir

    def extract_docx(self):
        """

        Returns
        -------
        DOCXParser for Docx files.
        extract text and write images in img_dir

        """
        output = docx2txt.process(self.file, self.img_dir)
        return output


class XLSXParser:

    def __init__(self, filename):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename

    def extract_xlsx(self):
        """

        Returns
        -------
        XLSXParser for XLSX and XLS files.
        """
        out_list = []
        book = xlrd.open_workbook(self.file)
        for val in range(len(book.sheet_names())):
            sheet = book.sheet_by_index(val)
            for res in range(sheet.nrows):
                output = " " + " ".join(str(val_) for val_ in (sheet.row_values(res)))
                out_list.append(output)
        return out_list


class CSVParser:

    def __init__(self, filename, delimiter):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        delimiter : By default ','. Can be changed if any other delimiter is needed.

        """
        self.file = filename
        self.delimiter = delimiter

    def extract_csv(self):
        """

        Returns
        -------
        CSVParser for csv files.

        """
        with open(self.file) as file:
            output = csv.reader(file, delimiter=self.delimiter)
            return ' '.join([' '.join(row) for row in output])


class ImageParser:

    def __init__(self, filename, tess_path=None):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        tess_path : The path to the tesseract cmd (Only for windows.)
        """
        self.file = filename
        self.path = tess_path

    def extract_image(self):
        """

        Returns
        -------
        ImageParser for Image formats.

        """
        out_list = []
        if self.path is not None:
            pytesseract.pytesseract.tesseract_cmd = self.path
            img = Image.open(self.file)
            text = pytesseract.image_to_string(img)
            out_list.append(text)
        else:
            img = Image.open(self.file)
            text = pytesseract.image_to_string(img)
            out_list.append(text)
        return out_list


class PDFParser:

    def __init__(self, filename, insights=False):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        insights : True by default. False if the dataframe is not needed.
        """
        self.file = filename
        self.insights = insights

    def extract_pypdf2(self):
        """

        Returns
        -------
        PDFParser for pdf files.

        """
        contents = []
        text = ' '
        pdfFileObj = open(self.file, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        pdfPages = pdfReader.getNumPages()
        if pdfPages == 1:
            for val in range(pdfReader.numPages):
                pageObject = pdfReader.getPage(val)
                text = text + pageObject.extractText()
            contents.append(text)
            if self.insights:
                conv = conversion(self.file)
                __conv = conv.convert()
                insights = generate_insights(__conv, df)
                __insights = insights.generate_df()
                remove_files(__conv)
                return __insights, contents
            else:
                return contents

        if pdfPages >= 2:
            pdf_multipage_df = pd.DataFrame()
            for val in range(pdfReader.numPages):
                pageObject = pdfReader.getPage(val)
                text = text + pageObject.extractText()
            contents.append(text)
            if self.insights:
                df_list = []
                pdffile = self.file
                os.mkdir("tempdir")
                tempdir = "tempdir/"
                if os.path.isdir(tempdir):
                    shutil.rmtree(tempdir)
                os.mkdir("tempdir/converted_files")
                images = convert_from_path(pdffile)
                converted_files = tempdir + "converted_files/"
                for val in range(len(images)):
                    images[val - 1].save(converted_files + str(val) + ".jpg", "JPEG")
                jpgfiles = os.listdir(converted_files)
                output_files = [converted_files + os.sep + _val for _val in jpgfiles if _val[-3:].upper() == "JPG"]
                for val in range(len(output_files)):
                    insights = generate_insights(output_files[val], df)
                    __insights = insights.generate_df()
                    page = [val] * len(__insights)
                    __insights["page_num"] = page
                    df_list.append(__insights)
                    pdf_multipage_df = pd.concat([pdf_multipage_df, __insights])
                shutil.rmtree(tempdir)
                df1 = pdf_multipage_df.reset_index()
                df1 = df1.drop("index", 1)
                return df1, contents
            else:
                return contents

    def extract_pdfplumber(self):
        """

        Returns
        -------
        Works as an alternative for PyPDF2.
        """
        out_list = []
        with pdfplumber.open(self.file) as pdf:
            for val in range(len(pdf.pages)):
                page = pdf.pages[val]
                output = page.extract_text()
                out_list.append(output)
        return out_list


class TXTParser:

    def __init__(self, filename):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename

    def extract_txt(self):
        """

        Returns
        -------
        TXTParser for txt, log or no extension files.
        """
        with open(self.file) as file:
            return file.read()


class PPTXParser:

    def __init__(self, filename):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename

    def extract_pptx(self):
        """

        Returns
        -------
        PPTXParser for pptx files.
        """
        text = []
        paper = Presentation(self.file)
        for slide in paper.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    stripped = paragraph.text.strip()
                    if stripped:
                        text.append(paragraph.text)
        return text


class speech_to_text:

    def __init__(self, filename):
        """

        Parameters
        ----------
        filename : The file that needs to be processed.
        """
        self.file = filename

    def extract_audio(self):
        """

        Returns
        -------
        speech_to_text for mp3, wav files.
        """
        output_audio = []
        os.mkdir("tempdir")
        dst_file = mp3_to_wav(self.file, "tempdir/sample.wav", format="wav")
        output = sr.AudioFile(dst_file)
        recog = sr.Recognizer()
        with output as source:
            audio = recog.record(source)
        output_audio.append(recog.recognize_google(audio))
        shutil.rmtree("tempdir")
        return output_audio
